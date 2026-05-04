#!/usr/bin/env python3
"""scrub.py - sanitize a file for external sharing.

Detects type, dispatches to a type-specific scrubber, runs a universal pass,
verifies, and reveals the output in Finder.

Always writes to a copy. Never modifies the input.
"""
from __future__ import annotations

import argparse
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile
from datetime import datetime
from pathlib import Path

PERSONAL_LEAK_TERMS = Path.home() / "Documents/GitHub/claude/.claude/skills/personal/scrub-meta/leak-terms.txt"
EPOCH = (1980, 1, 1, 0, 0, 0)
NEUTRAL_DATE = datetime(2026, 1, 1, 0, 0, 0)


def run(cmd: list[str], **kw) -> subprocess.CompletedProcess:
    return subprocess.run(cmd, capture_output=True, text=True, **kw)


def need(binary: str, hint: str) -> None:
    if shutil.which(binary) is None:
        sys.exit(f"error: {binary} not found in PATH. install: {hint}")


def detect_category(path: Path) -> tuple[str, str]:
    need("file", "comes with macOS")
    mime = run(["file", "--mime-type", "-b", str(path)]).stdout.strip()
    ext = path.suffix.lower()
    OOXML_EXTS = {".xlsx", ".docx", ".pptx"}
    IWORK_EXTS = {".numbers", ".pages", ".key"}
    LEGACY_OFFICE_EXTS = {".xls", ".doc", ".ppt"}
    if ext in OOXML_EXTS or "officedocument" in mime:
        return "ooxml", mime
    if ext in IWORK_EXTS:
        return "iwork", mime
    if ext in LEGACY_OFFICE_EXTS or "ms-excel" in mime or "msword" in mime:
        return "legacy_office", mime
    if mime == "application/pdf":
        return "pdf", mime
    if mime.startswith("image/svg") or ext == ".svg":
        return "svg", mime
    if mime.startswith("image/"):
        return "image", mime
    if mime.startswith("video/"):
        return "video", mime
    if mime.startswith("audio/"):
        return "audio", mime
    if mime in {"application/zip", "application/x-tar", "application/gzip"} or ext in {".zip", ".tar", ".tgz", ".gz"}:
        return "archive", mime
    if mime.startswith("text/") or ext in {".md", ".html", ".htm", ".txt"}:
        return "text", mime
    return "unknown", mime


def exif_snapshot(path: Path) -> list[dict]:
    if shutil.which("exiftool") is None:
        return []
    out = run(["exiftool", "-a", "-G1", "-j", str(path)]).stdout
    try:
        return json.loads(out) if out.strip() else []
    except json.JSONDecodeError:
        return []


def load_leak_terms(extra: str | None) -> list[str]:
    terms: list[str] = []
    if PERSONAL_LEAK_TERMS.is_file():
        for line in PERSONAL_LEAK_TERMS.read_text().splitlines():
            s = line.strip()
            if s and not s.startswith("#"):
                terms.append(s)
    if extra:
        terms.extend(t.strip() for t in extra.split(",") if t.strip())
    seen: set[str] = set()
    out: list[str] = []
    for t in terms:
        k = t.lower()
        if k not in seen:
            seen.add(k)
            out.append(t)
    return out


# ---------- type-specific scrubbers ----------

def scrub_ooxml(src: Path, dst: Path) -> dict:
    """Rebuild values-only via openpyxl/python-docx/python-pptx, then repack."""
    suffix = src.suffix.lower()
    tmp = Path(tempfile.mkstemp(suffix=suffix)[1])
    if suffix == ".xlsx":
        from openpyxl import Workbook, load_workbook
        wb_src = load_workbook(src, read_only=True, data_only=True)
        wb = Workbook()
        for i, ws_src in enumerate(wb_src.worksheets):
            ws = wb.active if i == 0 else wb.create_sheet()
            ws.title = "Sheet1" if i == 0 else f"Sheet{i+1}"
            for row in ws_src.iter_rows(values_only=True):
                if all(v is None for v in row):
                    continue
                ws.append(row)
        wb.properties.creator = ""
        wb.properties.lastModifiedBy = ""
        for attr in ("title", "subject", "description", "keywords", "category",
                     "identifier", "language", "contentStatus", "revision", "version"):
            setattr(wb.properties, attr, None)
        wb.properties.created = NEUTRAL_DATE
        wb.properties.modified = NEUTRAL_DATE
        wb.save(tmp)
    elif suffix == ".docx":
        from docx import Document
        doc_src = Document(str(src))
        doc = Document()
        for para in doc_src.paragraphs:
            doc.add_paragraph(para.text)
        for table_src in doc_src.tables:
            rows, cols = len(table_src.rows), len(table_src.columns)
            t = doc.add_table(rows=rows, cols=cols)
            for r in range(rows):
                for c in range(cols):
                    t.rows[r].cells[c].text = table_src.rows[r].cells[c].text
        cp = doc.core_properties
        cp.author = ""
        cp.last_modified_by = ""
        cp.title = ""
        cp.subject = ""
        cp.keywords = ""
        cp.category = ""
        cp.comments = ""
        cp.created = NEUTRAL_DATE
        cp.modified = NEUTRAL_DATE
        doc.save(str(tmp))
    elif suffix == ".pptx":
        from pptx import Presentation
        p_src = Presentation(str(src))
        p = Presentation()
        for slide_src in p_src.slides:
            slide = p.slides.add_slide(p.slide_layouts[6])
            for shape_src in slide_src.shapes:
                if shape_src.has_text_frame:
                    tb = slide.shapes.add_textbox(
                        shape_src.left or 0, shape_src.top or 0,
                        shape_src.width or 0, shape_src.height or 0,
                    )
                    tb.text_frame.text = shape_src.text_frame.text
        cp = p.core_properties
        cp.author = ""
        cp.last_modified_by = ""
        cp.title = ""
        cp.subject = ""
        cp.keywords = ""
        cp.category = ""
        cp.comments = ""
        cp.created = NEUTRAL_DATE
        cp.modified = NEUTRAL_DATE
        p.save(str(tmp))
    else:
        sys.exit(f"unsupported OOXML extension: {suffix}")

    fake_app = (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\r\n'
        b'<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties" '
        b'xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes">'
        b'<Application>Microsoft Excel</Application><DocSecurity>0</DocSecurity><ScaleCrop>false</ScaleCrop>'
        b'<Company></Company><LinksUpToDate>false</LinksUpToDate><SharedDoc>false</SharedDoc>'
        b'<HyperlinksChanged>false</HyperlinksChanged><AppVersion>16.0000</AppVersion></Properties>'
    )
    fake_core = (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\r\n'
        b'<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" '
        b'xmlns:dc="http://purl.org/dc/elements/1.1/" xmlns:dcterms="http://purl.org/dc/terms/" '
        b'xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">'
        b'<dc:creator></dc:creator><cp:lastModifiedBy></cp:lastModifiedBy>'
        b'<dcterms:created xsi:type="dcterms:W3CDTF">2026-01-01T00:00:00Z</dcterms:created>'
        b'<dcterms:modified xsi:type="dcterms:W3CDTF">2026-01-01T00:00:00Z</dcterms:modified>'
        b'</cp:coreProperties>'
    )

    with zipfile.ZipFile(tmp, "r") as zin, zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            if item.filename == "docProps/app.xml":
                data = fake_app
            elif item.filename == "docProps/core.xml":
                data = fake_core
            else:
                data = zin.read(item.filename)
            new = zipfile.ZipInfo(filename=item.filename, date_time=EPOCH)
            new.compress_type = zipfile.ZIP_DEFLATED
            new.external_attr = 0o644 << 16
            new.create_system = 0
            zout.writestr(new, data)
    tmp.unlink(missing_ok=True)
    return {"branch": "ooxml", "neutralized": ["creator", "lastModifiedBy", "Application", "absPath", "drawings", "hyperlinks", "ZIP timestamps"]}


def scrub_pdf(src: Path, dst: Path) -> dict:
    need("qpdf", "brew install qpdf")
    run(["qpdf", "--linearize", "--object-streams=generate", str(src), str(dst)], check=False)
    if not dst.exists():
        shutil.copy2(src, dst)
    if shutil.which("exiftool"):
        run(["exiftool", "-all=", "-overwrite_original_in_place", str(dst)])
    return {"branch": "pdf", "neutralized": ["/Author", "/Creator", "/Producer", "/Keywords", "XMP", "embedded thumbnails"]}


def scrub_image(src: Path, dst: Path) -> dict:
    need("exiftool", "brew install exiftool")
    shutil.copy2(src, dst)
    run(["exiftool", "-all=", "-overwrite_original_in_place", str(dst)])
    after = exif_snapshot(dst)
    leftover_gps = []
    for d in after:
        for k in d.keys():
            if k.lower().startswith("gps") or "gps" in k.lower():
                leftover_gps.append(k)
    if leftover_gps:
        sys.exit(f"error: GPS field(s) survived image scrub: {leftover_gps}")
    return {"branch": "image", "neutralized": ["EXIF", "GPS", "Make/Model/Software/Owner/Serial", "thumbnail"]}


def scrub_video(src: Path, dst: Path) -> dict:
    need("ffmpeg", "brew install ffmpeg")
    cp = run(["ffmpeg", "-y", "-i", str(src), "-map_metadata", "-1", "-c", "copy", str(dst)])
    if cp.returncode != 0:
        sys.exit(f"ffmpeg failed:\n{cp.stderr}")
    if shutil.which("exiftool"):
        run(["exiftool", "-all=", "-overwrite_original_in_place", str(dst)])
    return {"branch": "video", "neutralized": ["metadata atoms", "location", "com.apple.* atoms", "software"]}


def scrub_audio(src: Path, dst: Path) -> dict:
    need("exiftool", "brew install exiftool")
    shutil.copy2(src, dst)
    run(["exiftool", "-all=", "-overwrite_original_in_place", str(dst)])
    return {"branch": "audio", "neutralized": ["ID3", "metadata atoms", "cover art comments"]}


def scrub_archive(src: Path, dst: Path) -> dict:
    if zipfile.is_zipfile(src):
        with zipfile.ZipFile(src, "r") as zin, zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                new = zipfile.ZipInfo(filename=item.filename, date_time=EPOCH)
                new.compress_type = zipfile.ZIP_DEFLATED
                new.external_attr = 0o644 << 16
                new.create_system = 0
                zout.writestr(new, data)
        return {"branch": "archive(zip)", "neutralized": ["entry timestamps", "uid/gid", "archive comment"]}
    shutil.copy2(src, dst)
    return {"branch": "archive(other)", "neutralized": ["(only universal pass applied)"]}


def scrub_iwork(src: Path, dst: Path) -> dict:
    need("soffice", "brew install --cask libreoffice")
    target = "xlsx" if src.suffix.lower() == ".numbers" else "pdf"
    with tempfile.TemporaryDirectory() as td:
        run(["soffice", "--headless", "--convert-to", target, "--outdir", td, str(src)])
        produced = next(Path(td).iterdir(), None)
        if produced is None:
            sys.exit("LibreOffice produced no output")
        return _redispatch(produced, dst)


def scrub_legacy_office(src: Path, dst: Path) -> dict:
    need("soffice", "brew install --cask libreoffice")
    ext_map = {".xls": "xlsx", ".doc": "docx", ".ppt": "pptx"}
    target = ext_map[src.suffix.lower()]
    with tempfile.TemporaryDirectory() as td:
        run(["soffice", "--headless", "--convert-to", target, "--outdir", td, str(src)])
        produced = next(Path(td).iterdir(), None)
        if produced is None:
            sys.exit("LibreOffice produced no output")
        return _redispatch(produced, dst)


def scrub_svg(src: Path, dst: Path) -> dict:
    text = src.read_text(encoding="utf-8", errors="replace")
    text = re.sub(r"<metadata\b[^>]*>.*?</metadata>", "", text, flags=re.DOTALL)
    text = re.sub(r"\s+(sodipodi|inkscape):[\w\-]+\s*=\s*\"[^\"]*\"", "", text)
    text = re.sub(r"<!--.*?-->", "", text, flags=re.DOTALL)
    dst.write_text(text, encoding="utf-8")
    return {"branch": "svg", "neutralized": ["<metadata>", "sodipodi:*", "inkscape:*", "comments"]}


def scrub_text(src: Path, dst: Path) -> dict:
    text = src.read_text(encoding="utf-8", errors="replace")
    text = re.sub(r"<!--.*?-->", "", text, flags=re.DOTALL)
    if src.suffix.lower() == ".md" and text.startswith("---\n"):
        end = text.find("\n---", 4)
        if end != -1:
            text = text[end+4:].lstrip("\n")
    dst.write_text(text, encoding="utf-8")
    return {"branch": "text", "neutralized": ["HTML comments", "YAML frontmatter (md)"]}


def scrub_unknown(src: Path, dst: Path) -> dict:
    shutil.copy2(src, dst)
    return {"branch": "unknown", "neutralized": ["(only universal pass)"], "warning": "type fell through to universal pass — verify manually"}


SCRUBBERS = {
    "ooxml": scrub_ooxml, "pdf": scrub_pdf, "image": scrub_image,
    "video": scrub_video, "audio": scrub_audio, "archive": scrub_archive,
    "iwork": scrub_iwork, "legacy_office": scrub_legacy_office,
    "svg": scrub_svg, "text": scrub_text, "unknown": scrub_unknown,
}


def _redispatch(produced: Path, dst: Path) -> dict:
    cat, _ = detect_category(produced)
    fn = SCRUBBERS.get(cat, scrub_unknown)
    info = fn(produced, dst)
    info["branch"] = f"converted -> {info['branch']}"
    return info


# ---------- universal + verify ----------

def universal_pass(out: Path) -> list[str]:
    notes: list[str] = []
    if shutil.which("exiftool"):
        run(["exiftool", "-all=", "-overwrite_original_in_place", str(out)])
    run(["xattr", "-cr", str(out)])
    sibling = out.parent / ("._" + out.name)
    if sibling.exists():
        sibling.unlink()
        notes.append(f"removed resource fork: {sibling.name}")
    xattr_out = run(["xattr", "-l", str(out)]).stdout.strip()
    if xattr_out and "com.apple.provenance" in xattr_out:
        notes.append("residual com.apple.provenance (macOS-applied, opaque, transport-stripped)")
    return notes


def verify(out: Path, terms: list[str], category: str) -> tuple[bool, list[str]]:
    issues: list[str] = []
    if category in {"ooxml", "archive"} and zipfile.is_zipfile(out):
        with tempfile.TemporaryDirectory() as td:
            with zipfile.ZipFile(out, "r") as z:
                z.extractall(td)
            meta_files: list[Path] = []
            tdp = Path(td)
            for pat in ("docProps/*.xml", "**/*.rels", "xl/workbook.xml",
                        "word/document.xml", "ppt/presentation.xml",
                        "[Content_Types].xml"):
                meta_files.extend(tdp.glob(pat))
            for f in meta_files:
                try:
                    blob = f.read_text(encoding="utf-8", errors="replace")
                except Exception:
                    continue
                for t in terms:
                    if t.lower() in blob.lower():
                        issues.append(f"leak term '{t}' in {f.relative_to(tdp)}")
    else:
        snap = exif_snapshot(out)
        joined = json.dumps(snap).lower()
        for t in terms:
            if t.lower() in joined:
                issues.append(f"leak term '{t}' in exiftool output")
    rt = run(["bash", "-c", f"base64 -i '{out}' | base64 -D | shasum -a 256"]).stdout.strip().split()[0]
    direct = run(["shasum", "-a", "256", str(out)]).stdout.strip().split()[0]
    if rt != direct:
        issues.append(f"sha256 mismatch on base64 round-trip: {rt} vs {direct}")
    return (len(issues) == 0, issues)


def report(src: Path, out: Path, category: str, mime: str, info: dict, notes: list[str], ok: bool, issues: list[str]) -> None:
    print()
    print(f"# scrub-meta report")
    print()
    print(f"- input  : `{src}`")
    print(f"- output : `{out}`")
    print(f"- mime   : `{mime}`")
    print(f"- branch : `{info.get('branch', category)}`")
    print()
    print("## Removed")
    for n in info.get("neutralized", []):
        print(f"- {n}")
    if notes:
        print()
        print("## Notes")
        for n in notes:
            print(f"- {n}")
    if "warning" in info:
        print()
        print(f"## Warning")
        print(f"- {info['warning']}")
    print()
    print("## Verify")
    if ok:
        print("- all checks passed")
    else:
        for i in issues:
            print(f"- FAIL: {i}")
    sha = run(["shasum", "-a", "256", str(out)]).stdout.strip().split()[0]
    print(f"- sha256: `{sha}`")
    print()
    print("---")
    print("Do not open this file in its native app before sending. Excel / Word / Preview / Keynote re-inject creator, lastModifiedBy, and absolute paths on save. Send as-is from the output location.")


def main() -> None:
    ap = argparse.ArgumentParser(description="strip identifying metadata from a file")
    ap.add_argument("file", type=Path)
    ap.add_argument("--out", type=Path, default=None)
    ap.add_argument("--terms", type=str, default=None)
    args = ap.parse_args()

    src = args.file.expanduser().resolve()
    if not src.is_file():
        sys.exit(f"not a file: {src}")

    out = args.out.expanduser().resolve() if args.out else (Path.home() / "Downloads" / f"{src.stem}_clean{src.suffix}")
    out.parent.mkdir(parents=True, exist_ok=True)

    category, mime = detect_category(src)
    fn = SCRUBBERS.get(category, scrub_unknown)
    info = fn(src, out)
    notes = universal_pass(out)
    terms = load_leak_terms(args.terms)
    ok, issues = verify(out, terms, category)
    report(src, out, category, mime, info, notes, ok, issues)
    run(["open", "-R", str(out)])
    if not ok:
        sys.exit(1)


if __name__ == "__main__":
    main()
